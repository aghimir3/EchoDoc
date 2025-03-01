import logging
import json
import os
import pandas as pd
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from app.services.synthetic_data import SyntheticDataGenerator

logger = logging.getLogger(__name__)

class DocumentParser:
    """
    Service for parsing various document types.
    
    Converts documents (PDF, DOCX, PPTX, CSV, Excel, Text) into text,
    splits them into chunks, and optionally generates synthetic JSONL data.
    """
    def __init__(self) -> None:
        """Initialize DocumentParser."""
        self.logger = logger

    async def process_file(self, local_file_path: str, file_type: str, job_id: int = None) -> dict:
        """
        Process a file and generate text artifacts.
        
        :param local_file_path: Local file path to the document.
        :param file_type: MIME type of the document.
        :param job_id: Optional job ID for logging purposes.
        :return: Dictionary with keys 'jsonl' (synthetic JSONL string) and 'chunks' (list of text chunks).
        :raises Exception: If any error occurs during file processing.
        """
        try:
            text = self._parse_file(local_file_path, file_type)
            chunks = self.chunk_text(text)
            generator = SyntheticDataGenerator()
            # Await the asynchronous generation of synthetic JSONL data.
            jsonl_str = await generator.generate_synthetic_jsonl(chunks, job_id=job_id) if chunks else None
            return {"jsonl": jsonl_str, "chunks": chunks}
        except Exception as e:
            self.logger.error(f"Error processing file {local_file_path}: {str(e)}", exc_info=True)
            raise

    def _parse_file(self, local_file_path: str, file_type: str) -> str:
        """
        Dispatch file parsing based on file type.
        
        :param local_file_path: Local file path.
        :param file_type: MIME type.
        :return: Extracted text.
        :raises ValueError: If file type is unsupported.
        """
        parsers = {
            "application/pdf": self.parse_pdf,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": self.parse_docx,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": self.parse_pptx,
            "text/plain": self.read_text_file,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": self.parse_dataframe,
            "text/csv": self.parse_dataframe,
        }
        if file_type in parsers:
            return parsers[file_type](local_file_path, file_type) if file_type in [
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "text/csv"
            ] else parsers[file_type](local_file_path)
        raise ValueError(f"Unsupported file type: {file_type}")

    def parse_pdf(self, file_path: str) -> str:
        """
        Extract text from a PDF file.
        
        :param file_path: Path to the PDF.
        :return: Extracted text.
        :raises ValueError: If parsing fails.
        """
        try:
            reader = PdfReader(file_path)
            text = "".join(page.extract_text() or "" for page in reader.pages)
            self.logger.debug(f"Extracted text from PDF: {len(text)} characters")
            return text
        except Exception as e:
            self.logger.error(f"Failed to parse PDF {file_path}: {str(e)}")
            raise ValueError(f"Failed to parse PDF: {str(e)}")

    def parse_docx(self, file_path: str) -> str:
        """
        Extract text from a DOCX file.
        
        :param file_path: Path to the DOCX.
        :return: Extracted text.
        :raises ValueError: If parsing fails.
        """
        try:
            doc = DocxDocument(file_path)
            text = "\n".join(para.text for para in doc.paragraphs)
            self.logger.debug(f"Extracted text from DOCX: {len(text)} characters")
            return text
        except Exception as e:
            self.logger.error(f"Failed to parse DOCX {file_path}: {str(e)}")
            raise ValueError(f"Failed to parse DOCX: {str(e)}")

    def parse_pptx(self, file_path: str) -> str:
        """
        Extract text from a PPTX file.
        
        :param file_path: Path to the PPTX.
        :return: Extracted text.
        :raises ValueError: If parsing fails.
        """
        try:
            prs = Presentation(file_path)
            text = "\n".join(
                shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
            )
            self.logger.debug(f"Extracted text from PPTX: {len(text)} characters")
            return text
        except Exception as e:
            self.logger.error(f"Failed to parse PPTX {file_path}: {str(e)}")
            raise ValueError(f"Failed to parse PPTX: {str(e)}")

    def read_text_file(self, file_path: str) -> str:
        """
        Read a plain text file.
        
        :param file_path: Path to the text file.
        :return: File content.
        :raises ValueError: If reading fails.
        """
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            self.logger.debug(f"Read text file {file_path}: {len(text)} characters")
            return text
        except Exception as e:
            self.logger.error(f"Failed to read text file {file_path}: {str(e)}")
            raise ValueError(f"Failed to read text file: {str(e)}")

    def parse_dataframe(self, file_path: str, file_type: str) -> str:
        """
        Parse a spreadsheet (Excel or CSV) and convert its content to CSV text.
        
        :param file_path: Path to the file.
        :param file_type: MIME type indicating Excel or CSV.
        :return: CSV string.
        :raises ValueError: If parsing fails.
        """
        try:
            if file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                df = pd.read_excel(file_path)
            elif file_type == "text/csv":
                df = pd.read_csv(file_path)
            else:
                raise ValueError(f"Unsupported dataframe file type: {file_type}")
            text = df.to_csv(index=False)
            self.logger.debug(f"Converted dataframe to text with {len(text)} characters")
            return text
        except Exception as e:
            self.logger.error(f"Failed to parse dataframe from {file_path}: {str(e)}")
            raise ValueError(f"Failed to parse dataframe: {str(e)}")

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000) -> list:
        """
        Split text into chunks of roughly `chunk_size` words.
        
        :param text: Input text.
        :param chunk_size: Number of words per chunk.
        :return: List of text chunks.
        :raises ValueError: If chunking fails.
        """
        try:
            words = text.split()
            chunks = [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size)]
            logger.debug(f"Chunked text into {len(chunks)} chunks")
            return chunks
        except Exception as e:
            logger.error(f"Failed to chunk text: {str(e)}")
            raise ValueError(f"Failed to chunk text: {str(e)}")
